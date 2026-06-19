"""Importador de documentos: convierte PDF, DOCX, TXT, HTML, MD, EPUB, XLSX, PPTX,
ODT, RTF y CSV a pasajes indexables.

Funciona completamente offline sin GPU. Extrae texto, normaliza y divide en fragmentos.
Todos los formatos tienen fallback stdlib para funcionar sin dependencias externas.
"""

import json
import os
import re
from pathlib import Path
from datetime import datetime

from core.nlu import tokenize, BASE_DIR, DATA_DIR
from core.document_cleaner import clean as clean_text, clean_stats

IMPORTS_DIR = os.path.join(BASE_DIR, "imports")
DOCUMENTS_DIR = os.path.join(BASE_DIR, "documents")
IMPORT_REGISTRY = os.path.join(DATA_DIR, "import_registry.json")


def _load_registry():
    """Carga el registro de documentos importados."""
    if os.path.exists(IMPORT_REGISTRY):
        with open(IMPORT_REGISTRY, encoding="utf-8") as f:
            return json.load(f)
    return {}


def _save_registry(registry):
    """Guarda el registro de documentos importados."""
    os.makedirs(os.path.dirname(IMPORT_REGISTRY), exist_ok=True)
    with open(IMPORT_REGISTRY, "w", encoding="utf-8") as f:
        json.dump(registry, f, indent=2, ensure_ascii=False)


def _extract_pdf_text(pdf_path):
    """Extrae texto de PDF. Usa pdfplumber si está disponible, sino lee binarios."""
    try:
        import pdfplumber
        with pdfplumber.open(pdf_path) as pdf:
            text = ""
            for page in pdf.pages:
                text += page.extract_text() or ""
                text += "\n"
            return text.strip()
    except ImportError:
        return _extract_pdf_text_fallback(pdf_path)


def _extract_pdf_text_fallback(pdf_path):
    """Fallback para PDF sin pdfplumber: intenta extraer texto bruto."""
    try:
        with open(pdf_path, "rb") as f:
            data = f.read()
        text = ""
        in_text = False
        buffer = []
        for byte in data:
            if byte >= 32 and byte < 127:
                buffer.append(chr(byte))
            elif byte == 10 or byte == 13:
                if buffer:
                    text += "".join(buffer) + "\n"
                    buffer = []
                in_text = False
        if buffer:
            text += "".join(buffer)
        cleaned = re.sub(r"[^\w\s\.,;:!?\-áéíóúñüÁÉÍÓÚÑÜ]", "", text)
        return cleaned.strip() or "[PDF: no se pudo extraer texto — requiere pdfplumber]"
    except Exception as e:
        return f"[Error leyendo PDF: {e}]"


def _extract_docx_text(docx_path):
    """Extrae texto de DOCX. Usa python-docx si está disponible, sino descomprime ZIP."""
    try:
        from docx import Document
        doc = Document(docx_path)
        text = "\n".join(para.text for para in doc.paragraphs)
        return text.strip()
    except ImportError:
        return _extract_docx_text_fallback(docx_path)


def _extract_docx_text_fallback(docx_path):
    """Fallback para DOCX sin python-docx: descomprime el ZIP y busca XML."""
    try:
        import zipfile
        import xml.etree.ElementTree as ET
        with zipfile.ZipFile(docx_path, "r") as zf:
            if "word/document.xml" not in zf.namelist():
                return "[DOCX: estructura no estándar]"
            xml_str = zf.read("word/document.xml")
        root = ET.fromstring(xml_str)
        ns = {"w": "http://schemas.openxmlformats.org/wordprocessingml/2006/main"}
        text_elements = root.findall(".//w:t", ns)
        text = "".join(elem.text for elem in text_elements if elem.text)
        return text.strip()
    except Exception as e:
        return f"[Error leyendo DOCX: {e}]"


def _extract_html_text(html_path):
    """Extrae texto de HTML. Usa BeautifulSoup si está disponible, sino regex."""
    try:
        from bs4 import BeautifulSoup
        with open(html_path, encoding="utf-8", errors="replace") as f:
            soup = BeautifulSoup(f, "html.parser")
        for script in soup(["script", "style"]):
            script.decompose()
        text = soup.get_text(separator="\n")
        return re.sub(r"\n\s*\n", "\n\n", text).strip()
    except ImportError:
        return _extract_html_text_fallback(html_path)


def _extract_html_text_fallback(html_path):
    """Fallback para HTML sin BeautifulSoup: usa regex."""
    with open(html_path, encoding="utf-8", errors="replace") as f:
        html = f.read()
    html = re.sub(r"<script[^>]*>.*?</script>", "", html, flags=re.DOTALL | re.IGNORECASE)
    html = re.sub(r"<style[^>]*>.*?</style>", "", html, flags=re.DOTALL | re.IGNORECASE)
    text = re.sub(r"<[^>]+>", "", html)
    text = re.sub(r"&nbsp;", " ", text)
    text = re.sub(r"&[a-z]+;", "", text, flags=re.IGNORECASE)
    return re.sub(r"\n\s*\n", "\n\n", text).strip()


# ---------------------------------------------------------------------------
# Extractores: XLSX / XLS
# ---------------------------------------------------------------------------

def _extract_xlsx_text(path):
    """Extrae texto de Excel. Usa openpyxl si está disponible, sino descomprime ZIP+XML."""
    try:
        import openpyxl
        wb = openpyxl.load_workbook(path, read_only=True, data_only=True)
        parts = []
        for sheet in wb.worksheets:
            parts.append(f"## Hoja: {sheet.title}")
            for row in sheet.iter_rows(values_only=True):
                cells = [str(c) for c in row if c is not None and str(c).strip()]
                if cells:
                    parts.append(" | ".join(cells))
        wb.close()
        return "\n".join(parts)
    except ImportError:
        return _extract_xlsx_fallback(path)


def _extract_xlsx_fallback(path):
    """Fallback: descomprime el ZIP y lee sharedStrings + sheets."""
    try:
        import zipfile
        import xml.etree.ElementTree as ET
        with zipfile.ZipFile(path, "r") as zf:
            names = zf.namelist()
            # Cadenas compartidas
            strings = []
            if "xl/sharedStrings.xml" in names:
                root = ET.fromstring(zf.read("xl/sharedStrings.xml"))
                for si in root.findall(".//{http://schemas.openxmlformats.org/spreadsheetml/2006/main}t"):
                    strings.append(si.text or "")
            # Leer primera hoja disponible
            sheet_files = sorted(n for n in names if re.match(r"xl/worksheets/sheet\d+\.xml", n))
            parts = []
            for sf in sheet_files[:5]:  # max 5 hojas
                root = ET.fromstring(zf.read(sf))
                ns = "http://schemas.openxmlformats.org/spreadsheetml/2006/main"
                for row in root.findall(f".//{{{ns}}}row"):
                    cells = []
                    for c in row.findall(f"{{{ns}}}c"):
                        t = c.get("t", "")
                        v = c.find(f"{{{ns}}}v")
                        if v is not None and v.text:
                            if t == "s" and v.text.isdigit():
                                idx = int(v.text)
                                cells.append(strings[idx] if idx < len(strings) else "")
                            else:
                                cells.append(v.text)
                    if cells:
                        parts.append(" | ".join(cells))
        return "\n".join(parts) or "[XLSX: no se pudo leer contenido]"
    except Exception as e:
        return f"[Error leyendo XLSX: {e}]"


def _extract_xls_text(path):
    """Extrae texto de XLS (formato antiguo). Usa xlrd si está disponible."""
    try:
        import xlrd
        wb = xlrd.open_workbook(path)
        parts = []
        for sheet in wb.sheets():
            parts.append(f"## Hoja: {sheet.name}")
            for row_idx in range(sheet.nrows):
                cells = [str(sheet.cell_value(row_idx, c))
                         for c in range(sheet.ncols)
                         if str(sheet.cell_value(row_idx, c)).strip()]
                if cells:
                    parts.append(" | ".join(cells))
        return "\n".join(parts)
    except ImportError:
        return "[XLS: requiere 'pip install xlrd' para leer archivos .xls antiguos]"
    except Exception as e:
        return f"[Error leyendo XLS: {e}]"


# ---------------------------------------------------------------------------
# Extractores: PPTX
# ---------------------------------------------------------------------------

def _extract_pptx_text(path):
    """Extrae texto de PowerPoint. Usa python-pptx si está disponible, sino ZIP+XML."""
    try:
        from pptx import Presentation
        prs = Presentation(path)
        parts = []
        for i, slide in enumerate(prs.slides, 1):
            slide_texts = []
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    slide_texts.append(shape.text.strip())
            if slide_texts:
                parts.append(f"## Diapositiva {i}")
                parts.extend(slide_texts)
        return "\n\n".join(parts)
    except ImportError:
        return _extract_pptx_fallback(path)


def _extract_pptx_fallback(path):
    """Fallback: descomprime PPTX como ZIP y extrae texto de slides XML."""
    try:
        import zipfile
        import xml.etree.ElementTree as ET
        ns = "http://schemas.openxmlformats.org/drawingml/2006/main"
        parts = []
        with zipfile.ZipFile(path, "r") as zf:
            slide_files = sorted(
                n for n in zf.namelist()
                if re.match(r"ppt/slides/slide\d+\.xml", n)
            )
            for i, sf in enumerate(slide_files, 1):
                root = ET.fromstring(zf.read(sf))
                texts = [t.text for t in root.findall(f".//{{{ns}}}t") if t.text]
                if texts:
                    parts.append(f"## Diapositiva {i}")
                    parts.append(" ".join(texts))
        return "\n\n".join(parts) or "[PPTX: no se encontró texto en las diapositivas]"
    except Exception as e:
        return f"[Error leyendo PPTX: {e}]"


# ---------------------------------------------------------------------------
# Extractores: ODT
# ---------------------------------------------------------------------------

def _extract_odt_text(path):
    """Extrae texto de ODT (OpenDocument Text). Fallback puro ZIP+XML."""
    try:
        import zipfile
        import xml.etree.ElementTree as ET
        with zipfile.ZipFile(path, "r") as zf:
            if "content.xml" not in zf.namelist():
                return "[ODT: estructura no estándar]"
            root = ET.fromstring(zf.read("content.xml"))
        ns = "urn:oasis:names:tc:opendocument:xmlns:text:1.0"
        parts = []
        for elem in root.iter(f"{{{ns}}}p"):
            text = "".join(elem.itertext()).strip()
            if text:
                parts.append(text)
        return "\n".join(parts) or "[ODT: no se encontró texto]"
    except Exception as e:
        return f"[Error leyendo ODT: {e}]"


# ---------------------------------------------------------------------------
# Extractores: RTF
# ---------------------------------------------------------------------------

_RTF_CONTROL_RE  = re.compile(r"\\[a-z*]+[-\d]* ?")
_RTF_GROUP_RE    = re.compile(r"[{}]")
_RTF_SPECIAL_RE  = re.compile(r"\\['`][0-9a-fA-F]{2}")
_RTF_UNICODE_RE  = re.compile(r"\\u(-?\d+)\??")


def _extract_rtf_text(path):
    """Extrae texto de RTF. Parser mínimo sin dependencias externas."""
    try:
        with open(path, "rb") as f:
            raw = f.read()
        # Intentar decodificar como latin-1 (codificación común en RTF)
        text = raw.decode("latin-1", errors="replace")

        # Eliminar grupos de metadatos {\*\comando ...}
        text = re.sub(r"\{\\\*[^}]*\}", "", text)
        # Caracteres Unicode RTF: \uNNNN
        def _rtf_unicode(m):
            code = int(m.group(1))
            if code < 0:
                code += 65536
            try:
                return chr(code)
            except (ValueError, OverflowError):
                return ""
        text = _RTF_UNICODE_RE.sub(_rtf_unicode, text)
        # Caracteres especiales RTF \\'xx
        def _rtf_hex(m):
            try:
                return bytes.fromhex(m.group(0)[-2:]).decode("latin-1")
            except Exception:
                return ""
        text = _RTF_SPECIAL_RE.sub(_rtf_hex, text)
        # Párrafos
        text = text.replace("\\par", "\n").replace("\\line", "\n")
        text = text.replace("\\tab", "\t")
        # Quitar comandos de control y llaves
        text = _RTF_CONTROL_RE.sub("", text)
        text = _RTF_GROUP_RE.sub("", text)
        return text.strip() or "[RTF: no se pudo extraer texto]"
    except Exception as e:
        return f"[Error leyendo RTF: {e}]"


# ---------------------------------------------------------------------------
# Extractores: CSV
# ---------------------------------------------------------------------------

def _extract_csv_text(path):
    """Extrae texto de CSV convirtiéndolo a formato legible por BM25."""
    import csv
    parts = []
    encodings = ["utf-8", "utf-8-sig", "latin-1", "cp1252"]
    for enc in encodings:
        try:
            with open(path, encoding=enc, newline="") as f:
                # Detectar delimitador
                sample = f.read(4096)
                f.seek(0)
                try:
                    dialect = csv.Sniffer().sniff(sample, delimiters=",;\t|")
                except csv.Error:
                    dialect = csv.excel
                reader = csv.reader(f, dialect)
                headers = None
                for i, row in enumerate(reader):
                    if not any(row):
                        continue
                    if i == 0:
                        headers = [c.strip() for c in row]
                        parts.append("Columnas: " + " | ".join(headers))
                        continue
                    if headers:
                        # Convertir fila a "campo: valor" para que BM25 lo entienda
                        pairs = [f"{h}: {v.strip()}" for h, v in zip(headers, row) if v.strip()]
                        if pairs:
                            parts.append(". ".join(pairs))
                    else:
                        parts.append(" | ".join(c.strip() for c in row if c.strip()))
                    if i > 5000:  # límite de filas para no sobrecargar
                        parts.append(f"[... {i} filas en total, mostrando primeras 5000]")
                        break
            return "\n".join(parts)
        except (UnicodeDecodeError, LookupError):
            continue
    return "[CSV: no se pudo leer el archivo]"


def _normalize_text(text):
    """Normaliza texto: limpia espacios, corrige saltos de línea."""
    text = re.sub(r"\r\n", "\n", text)
    text = re.sub(r" +", " ", text)
    text = re.sub(r"\n{3,}", "\n\n", text)
    return text.strip()


def _split_into_chunks(text, max_words=300, overlap=50):
    """Divide texto en fragmentos para indexación (sin cortar palabras)."""
    sentences = re.split(r"(?<=[.!?])\s+", text)
    chunks = []
    current = []
    current_words = 0

    for sent in sentences:
        sent_words = len(sent.split())
        if current_words + sent_words > max_words and current:
            chunk_text = " ".join(current)
            chunks.append(chunk_text)
            current = current[-overlap:] if overlap < len(current) else current
            current_words = sum(len(s.split()) for s in current)
        current.append(sent)
        current_words += sent_words

    if current:
        chunks.append(" ".join(current))

    return [c.strip() for c in chunks if len(c.split()) > 10]


class DocumentImporter:
    """Importador de documentos multiformato."""

    SUPPORTED_FORMATS = (
        ".pdf", ".docx", ".txt", ".html", ".md",
        ".xlsx", ".xls", ".pptx", ".odt", ".rtf", ".csv",
    )

    def __init__(self):
        os.makedirs(IMPORTS_DIR, exist_ok=True)
        os.makedirs(DOCUMENTS_DIR, exist_ok=True)
        self.registry = _load_registry()

    def import_file(self, source_path, target_category=None, document_classifier=None):
        """
        Importa un documento: extrae texto, crea estructura, devuelve fragmentos.

        Args:
            source_path: ruta del archivo a importar
            target_category: categoría destino (si None, se auto-clasifica)
            document_classifier: instancia de DocumentClassifier para auto-clasificación

        Returns:
            dict con metadatos y fragmentos, o None si error
        """
        source_path = str(source_path)
        if not os.path.exists(source_path):
            return None

        ext = Path(source_path).suffix.lower()
        if ext not in self.SUPPORTED_FORMATS:
            return None

        filename = Path(source_path).name
        file_hash = self._file_hash(source_path)

        if file_hash in self.registry:
            return {"status": "already_imported", "file": filename}

        try:
            text = self._extract_text(source_path, ext)
            if not text or len(text) < 50:
                return {"status": "empty_or_unreadable", "file": filename}

            # Limpieza profunda: artefactos PDF, duplicados, ruido tipográfico
            raw_text = _normalize_text(text)
            # PDFs suelen necesitar limpieza agresiva; TXT/MD son más limpios
            aggressive = Path(source_path).suffix.lower() == ".pdf"
            text = clean_text(raw_text, aggressive=aggressive)
            stats = clean_stats(raw_text, text)

            chunks = _split_into_chunks(text, max_words=300, overlap=50)

            if target_category is None and document_classifier:
                target_category = document_classifier.classify(text)

            target_category = target_category or "general"

            result = {
                "filename": filename,
                "original_path": source_path,
                "category": target_category,
                "imported_at": datetime.now().isoformat(),
                "file_hash": file_hash,
                "text_length": len(text),
                "chunk_count": len(chunks),
                "chunks": chunks,
                "cleaning": stats,
            }

            self.registry[file_hash] = {
                "filename": filename,
                "category": target_category,
                "imported_at": result["imported_at"],
                "chunk_count": len(chunks),
            }
            _save_registry(self.registry)

            return result

        except Exception as e:
            return {"status": "error", "file": filename, "error": str(e)}

    def _extract_text(self, path, ext):
        """Extrae texto según el formato."""
        if ext == ".pdf":
            return _extract_pdf_text(path)
        elif ext == ".docx":
            return _extract_docx_text(path)
        elif ext == ".html":
            return _extract_html_text(path)
        elif ext == ".xlsx":
            return _extract_xlsx_text(path)
        elif ext == ".xls":
            return _extract_xls_text(path)
        elif ext == ".pptx":
            return _extract_pptx_text(path)
        elif ext == ".odt":
            return _extract_odt_text(path)
        elif ext == ".rtf":
            return _extract_rtf_text(path)
        elif ext == ".csv":
            return _extract_csv_text(path)
        else:
            # .txt, .md y cualquier texto plano
            for enc in ("utf-8", "utf-8-sig", "latin-1", "cp1252"):
                try:
                    with open(path, encoding=enc) as f:
                        return f.read()
                except UnicodeDecodeError:
                    continue
            with open(path, encoding="utf-8", errors="replace") as f:
                return f.read()

    def _file_hash(self, path):
        """Calcula un hash simple del archivo (mtime + size)."""
        try:
            stat = os.stat(path)
            return f"{stat.st_mtime}_{stat.st_size}"
        except OSError:
            return ""

    def get_import_status(self):
        """Devuelve estadísticas de importación."""
        return {
            "total_imported": len(self.registry),
            "registry": self.registry,
        }
